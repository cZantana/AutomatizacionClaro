from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import zipfile
import xml.etree.ElementTree as ET
from colorama import Back, Style
import pytesseract
from PIL import Image
import io
import os

import pytesseract

# Configura la ruta de Tesseract según tu sistema operativo
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'


# Función para obtener el nombre del tipo de objeto
def get_shape_type_name(shape_type):
    shape_type_dict = {
        MSO_SHAPE_TYPE.PICTURE: "Imagen",
        MSO_SHAPE_TYPE.AUTO_SHAPE: "Autoforma",
        MSO_SHAPE_TYPE.TEXT_BOX: "Cuadro de texto",
        MSO_SHAPE_TYPE.MEDIA: "Objeto multimedia",
        MSO_SHAPE_TYPE.FREEFORM: "Forma libre",
        MSO_SHAPE_TYPE.GROUP: "Grupo de objetos",
        MSO_SHAPE_TYPE.LINE: "Línea",
        MSO_SHAPE_TYPE.TABLE: "Tabla",
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: "Objeto OLE incrustado"
    }
    return shape_type_dict.get(shape_type, "Tipo desconocido")

# Función para extraer el texto de un Cuadro de texto
# Función para extraer texto de un cuadro de texto o una autoforma
def extract_text_from_textbox(shape):
    if shape.has_text_frame:
        text = shape.text_frame.text
        return text.strip()
    return "Sin texto"


# Función actualizada para desglosar los objetos dentro de un grupo (incluye posición)
# Función para desglosar los objetos dentro de un grupo (incluye posición)
def extract_group_shapes(group_shape):
    group_objects = []
    for shape in group_shape.shapes:
        shape_type_name = get_shape_type_name(shape.shape_type)
        group_objects.append({
            'id': shape.shape_id,
            'name': shape.name,
            'type': shape_type_name,
            'shape': shape,
            'position': (shape.left, shape.top)
        })
    return group_objects



# Función para extraer texto de una imagen
def extract_text_from_image(shape, pptx_path):
    # Extraer la relación de la imagen para obtener el archivo binario
    image_rel = None
    for rel in shape.part.rels.values():
        if "image" in rel.reltype:
            image_rel = rel
            break

    if image_rel:
        # Normalizar la ruta para eliminar '../' si existe
        image_path = os.path.normpath(f'ppt/{image_rel.target_ref}').replace('\\', '/')

        # Asegurar que la ruta esté dentro de 'ppt/media/'
        if not image_path.startswith('ppt/media/'):
            image_path = f'ppt/media/{os.path.basename(image_rel.target_ref)}'

        try:
            with zipfile.ZipFile(pptx_path, 'r') as pptx_zip:
                if image_path in pptx_zip.namelist():
                    with pptx_zip.open(image_path) as image_file:
                        image_data = image_file.read()
                        image = Image.open(io.BytesIO(image_data))

                        # Usar pytesseract para extraer texto de la imagen
                        extracted_text = pytesseract.image_to_string(image, lang='eng')
                        return extracted_text.strip() if extracted_text.strip() else "Sin texto detectado en la imagen"
                else:
                    return f"Archivo de imagen no encontrado en el archivo PPTX: {image_path}"
        except Exception as e:
            return f"Error al procesar la imagen: {e}"
    else:
        return "No se encontró la relación de la imagen."


# Función para procesar un grupo de objetos y extraer el texto
# Función para procesar un grupo de objetos y extraer el texto de cuadros de texto y autoformas
def process_group_shape(group_shape, pptx_path):
    """
    Procesa un grupo de objetos, desglosa cada uno y extrae el texto.
    Los objetos se ordenan de arriba a la izquierda hacia abajo a la derecha.
    """
    # Desglosar objetos dentro del grupo (reutilizando extract_group_shapes)
    objects_in_group = extract_group_shapes(group_shape)

    # Ordenar los objetos por posición: primero arriba (top), luego a la izquierda (left)
    objects_in_group.sort(key=lambda obj: (obj['position'][1], obj['position'][0]))

    extracted_texts = []  # Lista para almacenar los textos extraídos

    # Procesar cada objeto del grupo según su tipo
    for obj in objects_in_group:
        shape = obj['shape']

        # Extraer texto si es un cuadro de texto o una autoforma con texto
        if obj['type'] in ["Cuadro de texto", "Autoforma"]:
            text = extract_text_from_textbox(shape)
            if text:
                extracted_texts.append(text)

        # Extraer texto de imágenes usando OCR
        elif obj['type'] == "Imagen":
            text = extract_text_from_image(shape, pptx_path)
            if text:
                extracted_texts.append(text)

    # Combinar todos los textos extraídos en orden de lectura
    combined_text = "\n".join(extracted_texts) if extracted_texts else "Sin texto en el grupo."

    return combined_text


# Función para extraer texto de una tabla
def extract_text_from_table(shape):
    """
    Extrae el texto de una tabla recorriendo todas sus celdas.
    Devuelve el texto de cada celda en el orden de lectura (fila por fila).
    """
    if shape.has_table:
        table = shape.table
        extracted_text = []

        # Recorrer las filas y columnas de la tabla
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                text = cell.text.strip() if cell.text else ''
                row_text.append(text)
            extracted_text.append(" | ".join(row_text))  # Separar celdas con '|'

        combined_text = "\n".join(extracted_text)
        return combined_text if combined_text else "Sin texto en la tabla."

    return "Sin texto en la tabla."

import re

# Función para segmentar comentarios por enlaces
def segment_comment_by_links(comment_text):
    """
    Segmenta un comentario en partes basadas en los enlaces encontrados.
    Devuelve una lista de segmentos, cada uno con el texto asociado al enlace.
    """
    # Patrón regex para identificar URLs que comienzan con https://
    pattern = r'(.*?)(https://[^\s]+)'  # Captura texto antes del enlace y el enlace

    matches = re.findall(pattern, comment_text, re.DOTALL)

    # Almacenar los segmentos extraídos
    segments = []

    # Procesar coincidencias
    for idx, (text_before_link, link) in enumerate(matches, start=1):
        segment = f"{idx} enlace del comentario:\n{text_before_link.strip()}\n{link.strip()}"
        segments.append(segment)

    return segments

# Lista para almacenar los datos de todos los comentarios segmentados
data = []


# Cargar la presentación
pptx_path = 'Dia 1 Febrero.pptx'
prs = Presentation(pptx_path)

# 1. Cargar los autores de comentarios
authors_dict = {}

with zipfile.ZipFile(pptx_path, 'r') as pptx_zip:
    if 'ppt/authors.xml' in pptx_zip.namelist():
        with pptx_zip.open('ppt/authors.xml') as authors_file:
            tree = ET.parse(authors_file)
            root = tree.getroot()

            for author in root.findall('.//author'):
                author_id = author.get('id')
                name = author.get('name')
                authors_dict[author_id] = name

# 2. Extraer objetos y comentarios por diapositiva
with zipfile.ZipFile(pptx_path, 'r') as pptx_zip:
    for index, slide in enumerate(prs.slides, start=1):

        rels_path = f'ppt/slides/_rels/slide{index}.xml.rels'

        if rels_path in pptx_zip.namelist():
            with pptx_zip.open(rels_path) as rels_file:
                rels_tree = ET.parse(rels_file)
                rels_root = rels_tree.getroot()

                namespace = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}

                for rel in rels_root.findall('r:Relationship', namespace):
                    rel_type = rel.get('Type')
                    target = rel.get('Target')

                    if 'comments' in rel_type:
                        print(f'\n{Back.GREEN}Diapositiva {index}:\n{Style.RESET_ALL}')

                        shapes_info = {}
                        for shape in slide.shapes:
                            shape_type = shape.shape_type
                            shape_id = shape.shape_id
                            shape_type_name = get_shape_type_name(shape_type)
                            shapes_info[str(shape_id)] = {
                                'type': shape_type_name,
                                'shape': shape
                            }
                            print(f'Tipo de objeto: {shape_type_name} (Código: {shape_type}, ID: {shape_id})')

                        comment_file = f'ppt/comments/{target.split("/")[-1]}'
                        print(f' - Buscando archivo de comentarios: {comment_file}')

                        if comment_file in pptx_zip.namelist():
                            with pptx_zip.open(comment_file) as comment_file_data:
                                comment_tree = ET.parse(comment_file_data)
                                comment_root = comment_tree.getroot()

                                for comment in comment_root.findall('.//{*}cm'):
                                    author_id = comment.get('authorId')
                                    comment_id = comment.get('id')
                                    created_date = comment.get('created')

                                    markers = [
                                        comment.find('.//{*}spMk'),
                                        comment.find('.//{*}picMk'),
                                        comment.find('.//{*}grpSpMk'),
                                        comment.find('.//{*}graphicFrameMk'),
                                        comment.find('.//{*}sldMk')
                                    ]

                                    shape_id = 'Sin ID'
                                    creation_id = 'Sin creationId'
                                    linked_shape = None

                                    for marker in markers:
                                        if marker is not None:
                                            shape_id = marker.get('id', 'Sin ID')
                                            creation_id = marker.get('creationId', 'Sin creationId')

                                            linked_shape = shapes_info.get(shape_id)
                                            break

                                    comment_text_elem = comment.find('.//{*}t')
                                    comment_text = comment_text_elem.text.strip() if comment_text_elem is not None else 'Sin texto'

                                    # Segmentación de comentarios por enlaces
                                    segments = segment_comment_by_links(comment_text)

                                    author_name = authors_dict.get(author_id, f'Autor ID {author_id}')

                                    print(f'{Back.BLUE} - Comentario ID {comment_id}:{Style.RESET_ALL}')
                                    print(f'   * Autor: {author_name}')
                                    print(f'   * Texto del comentario: {comment_text}')
                                    # Imprimir los segmentos encontrados

                                    print(f'{Back.MAGENTA}   * Vinculado a Objeto ID: {shape_id}, creationId: {creation_id}{Style.RESET_ALL}')

                                    if linked_shape:
                                        print(f'   * Tipo de Objeto: {linked_shape["type"]}')

                                        # Procesamiento para Cuadro de texto
                                        if linked_shape["type"] == "Cuadro de texto":
                                            extracted_text = extract_text_from_textbox(linked_shape['shape'])
                                            print(f'{Back.YELLOW}   * Texto extraído del Cuadro de texto: {extracted_text}{Style.RESET_ALL}')

                                        if linked_shape["type"] == "Autoforma":
                                            extracted_text = extract_text_from_textbox(linked_shape['shape'])
                                            print(f'{Back.YELLOW}   * Texto extraído de la autoforma: {extracted_text}{Style.RESET_ALL}')

                                        if linked_shape["type"] == "Imagen":
                                            extracted_text = extract_text_from_image(linked_shape['shape'], pptx_path)
                                            print(f'{Back.YELLOW}   * Texto extraído de la Imagen: {extracted_text}{Style.RESET_ALL}')

                                        if linked_shape["type"] == "Grupo de objetos":
                                            extracted_text = process_group_shape(linked_shape['shape'], pptx_path)
                                            print(f'{Back.YELLOW}   * Texto extraído del Grupo: {extracted_text}{Style.RESET_ALL}')

                                        if linked_shape["type"] == "Tabla":
                                            extracted_text = extract_text_from_table(linked_shape['shape'])
                                            print(f'{Back.YELLOW}   * Texto extraído de la Tabla:\n{extracted_text}{Style.RESET_ALL}')

                                        for segment in segments:
                                            print(f'{Back.CYAN}{segment}{Style.RESET_ALL}')
                                            policy = segment.split('\n')[1] if '\n' in segment else ''
                                            link = segment.split('\n')[-1] if 'https://' in segment else ''

                                            data.append({
                                                'Número de diapositiva': index,
                                                'Política': policy,
                                                'Enlace Conectados': link,
                                                'Tipo de objeto vinculado al comentario': linked_shape["type"] if linked_shape else 'Sin tipo',
                                                'Texto Extraído del objeto': extracted_text if linked_shape else 'Sin texto'
                                            })

                                    print(f'   * Fecha de creación: {created_date}')
                        else:
                            print(' - Archivo de comentarios no encontrado.')


# Crear un DataFrame y exportar a Excel
import pandas as pd

df = pd.DataFrame(data)
output_path = 'Comentarios_Presentación.xlsx'
df.to_excel(output_path, index=False)

print(f'Archivo Excel generado: {output_path}')